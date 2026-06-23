import os
import sys
import re
from pptx import Presentation

def main():
    if len(sys.argv) < 2:
        print("Usage: python process_ppt.py <ppt_path> [output_pdf_path] [output_notes_path]")
        sys.exit(1)
        
    ppt_path = sys.argv[1]
    
    base = os.path.splitext(ppt_path)[0]
    pdf_path = sys.argv[2] if len(sys.argv) > 2 else base + ".pdf"
    notes_path = sys.argv[3] if len(sys.argv) > 3 else base + "_notes.txt"
    
    # 1. Extract notes from PPTX
    print(f"Opening PPTX file: {ppt_path}")
    prs = Presentation(ppt_path)
    notes_list = []
    
    for i, slide in enumerate(prs.slides):
        note = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            note = slide.notes_slide.notes_text_frame.text.strip()
        
        # If the slide notes are empty, we should output a default placeholder
        if not note:
            note = "本页无配音。"
        notes_list.append(note)
        print(f"Slide {i+1} Note: {note[:40]}...")
        
    # Write notes
    with open(notes_path, "w", encoding="utf-8") as f:
        f.write("\n---\n".join(notes_list))
    print(f"Successfully extracted notes to: {notes_path}")
    
    # 2. Export to PDF via PowerPoint COM Interop
    print("Exporting PPTX to PDF via PowerPoint Application...")
    ppt_abs = os.path.abspath(ppt_path)
    pdf_abs = os.path.abspath(pdf_path)
    
    success = False
    
    # Try win32com first
    try:
        import win32com.client
        print("Using win32com.client...")
        powerpoint = win32com.client.DispatchEx("PowerPoint.Application")
        # Visibile=True is safer to prevent hang on some machines
        powerpoint.Visible = True
        deck = powerpoint.Presentations.Open(ppt_abs, WithWindow=False)
        deck.SaveAs(pdf_abs, 32) # 32 represents ppSaveAsPDF
        deck.Close()
        powerpoint.Quit()
        print(f"Successfully exported PDF to: {pdf_abs}")
        success = True
    except Exception as e:
        print(f"win32com export failed: {e}")
        
    if not success:
        # Try comtypes as fallback
        try:
            import comtypes.client
            print("Using comtypes.client...")
            powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
            powerpoint.Visible = True
            deck = powerpoint.Presentations.Open(ppt_abs, WithWindow=False)
            deck.SaveAs(pdf_abs, 32)
            deck.Close()
            powerpoint.Quit()
            print(f"Successfully exported PDF via comtypes to: {pdf_abs}")
            success = True
        except Exception as e:
            print(f"comtypes export failed: {e}")
            
    if not success:
        print("ERROR: Could not export PPTX to PDF via PowerPoint COM interface.")
        print("Please ensure MS PowerPoint is installed and registered in COM.")
        sys.exit(1)
        
    print("PPT Processing Done!")

if __name__ == "__main__":
    main()
